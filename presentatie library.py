from pptx import Presentation
import os

from GUI import slidelijsten  # Zorg dat GUI.py 'slidelijsten' exporteert (globale variabele)
print(slidelijsten)

prs = Presentation()

# Maak een titelslide (optioneel)
def titelslideophalen():
    slide_layout = prs.slide_layouts[0]  # Titel-layout
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = slidelijsten[0][0]
    subtitle.text = "\n".join(slidelijsten[0][1:]).strip()

# Maak de inhoudslides
def inhoudslides():
    for nummer, zinnen in slidelijsten.items():  # 
        if nummer == 0: #sla de gegevens van de titelslide over
            continue

        slide_layout = prs.slide_layouts[1]  # De layout (van de library) voor "Titel en inhoud"
        slide = prs.slides.add_slide(slide_layout)

        title = slide.shapes.title
        content = slide.placeholders[1]

        

        # Dit is voor de eerste zin in de list
        if len(zinnen) > 0:
            title.text = zinnen[0].strip()

        # dit is voor de rest van de zinnen in de list
        if len(zinnen) > 1:
            inhoud = "\n".join(zinnen[1:]).strip()
            content.text = inhoud

def bestandopslaan():
    downloads = os.path.join(os.path.expanduser("~"), "Downloads")
    PPTXbestand = os.path.join(downloads, "Menu.pptx")

    prs.save(PPTXbestand)
    print(f"Presentatie opgeslagen als: {PPTXbestand}")

titelslideophalen()
inhoudslides()
bestandopslaan()