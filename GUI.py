import pygame
import pygame_gui
import re
from pptx.util import Pt
import tkinter as tk
from tkinter import filedialog
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
import re
#from uitlezen import bestandlezen

import os

global textbox, titel_error_weergeven

achtergrondkleur = []

pygame.init()

pygame.display.set_caption('Text to slide')
window_surface = pygame.display.set_mode((500, 500))

background = pygame.Surface((800, 600))
background.fill(pygame.Color("#FFFFFF"))

manager = pygame_gui.UIManager((500, 500))

manager.preload_fonts([
    {'name': 'noto_sans', 'point_size': 14, 'style': 'bold', 'antialiased': '1'},
    {'name': 'noto_sans', 'point_size': 18, 'style': 'bold', 'antialiased': '1'},
    {'name': 'noto_sans', 'point_size': 18, 'style': 'italic', 'antialiased': '1'},
    {'name': 'noto_sans', 'point_size': 18, 'style': 'regular', 'antialiased': '1'},
    {'name': 'noto_sans', 'point_size': 48, 'style': 'regular', 'antialiased': '1'},
    {'name': 'noto_sans', 'point_size': 18, 'style': 'bold_italic', 'antialiased': '1'}
])

importeerknop = pygame_gui.elements.UIButton(
    relative_rect=pygame.Rect((0, 425), (500, 75)),
    text='Import tekstbestand',
    manager=manager,
)

bestandnaamvak = pygame_gui.elements.UITextEntryLine(
    relative_rect=pygame.Rect((0, 375), (500, 50)),
    placeholder_text="Typ hier de naam van je PPTX",
    manager=manager
)

textbox = pygame_gui.elements.UITextBox(
    relative_rect=pygame.Rect((0, 75), (500, 300)),
    html_text="<font size=5>Upload een .txt bestand die je wil gebruiken om een .PPTX bestand te maken.</font>",
    manager=manager
)

titel_error_weergeven, kleur_error_woordgeving_weergeven, kleur_gekozen_error_weergeven, kleur_gekozen = False, False, False, False
def add_to_log(log_tekst):
    huidige_tekst = textbox.html_text
    nieuwe_tekst = huidige_tekst + f"<br>{log_tekst}"
    textbox.set_text(nieuwe_tekst)
    
    textbox.rebuild()
    textbox.scroll_position = 0.5 # automatisch naar onderkant scrollbar scrollen


titlebox = pygame_gui.elements.UITextBox(
    relative_rect=pygame.Rect((0, 0), (500, 75)),
    html_text='<font bold><font size=7>         Text To Slide</font></font>',
    manager=manager
)

# Kleur voor importeerknop
importeerknop.colours['normal_bg'] = pygame.Color("#000000")
importeerknop.colours['hovered_bg'] = pygame.Color("#232323")
importeerknop.colours['active_bg'] = pygame.Color("#000000")
importeerknop.rebuild()

clock = pygame.time.Clock()
is_running = True

while is_running:
    time_delta = clock.tick(60) / 1000.0
    for event in pygame.event.get():
        if event.type == pygame.QUIT:
            is_running = False

        if event.type == pygame_gui.UI_BUTTON_PRESSED:

            if event.ui_element == importeerknop   :
                import tkinter as tk
                from tkinter import filedialog

                root = tk.Tk()
                root.withdraw()  
                bestand = filedialog.askopenfilename(
                    title="Voeg tekstbestand toe:",
                    filetypes=[("*Tekstbestand", "*.txt")]
                )

                if bestand:
                    def bestandlezen(bestand):
                        global titel_error_weergeven, kleur_error_woordgeving_weergeven, kleur_gekozen_error_weergeven, kleur_gekozen
                        afbeeldingen_per_slide = {}
                        #achtergrondkleur[] = RGBColor(255, 255, 255)

                        afbeeling_in_slide = True

                        with open(bestand, "r") as f:
                            text = f.read()
                        allezinnen = re.split(r'(?<!\\)\.', text) #laat zinnen met een backslash voor de punt in tact
                        allezinnen = [zin.replace('\.', '.') for zin in allezinnen] #veranderd de zinnen met \. erin in . zodat het correct weergeven wordt

                        slidelijsten = {}
                        slidenummer = 0

                        for zin in allezinnen:
                            zin = zin.strip()
                            zin = zin.replace('\\.', '.').replace('\\@', '@').replace('\\#', '#')


                            print(zin)
                            if zin.startswith("\\"):
                                if zin.startswith("\\@"):
                                    print("Nu doet hij het wel")
                                if zin.startswith("\\#"):
                                    # overslaan
                                    print("Hij slaat over")
                                else:
                                    if slidenummer == 0:
                                        slidenummer = 1
                                    if slidenummer not in slidelijsten:
                                        slidelijsten[slidenummer] = []
                                    zin = zin[1:]  # haalt \ weg uit presentatie
                                    slidelijsten[slidenummer].append(zin)

                            elif zin.startswith("#"):
                                zin = zin[1:]
                                if 0 in slidelijsten and slidelijsten[0] is not None:
                                    #de voorpagina mag maar twee zinnen krijgen dus als die meer dan dat krijgt wordt dat niet toegevoegd
                                    #en wordt er een error gegeven. en dat maar een keer vanwege de titel_error_weergeven variabel 
                                    if len(slidelijsten[0]) < 2:
                                        slidelijsten[0].append(zin)
                                    elif len(slidelijsten[0]) >= 2:
                                        if titel_error_weergeven == False:
                                            log_tekst = (
                                                f"<font size=5><b><font color='#ff0000'>Error: kon de volgende zin niet in presentatie zetten:</font></b> <i>'{zin}'</i> </font>"
                                                f"<font size=5><b><font color='#11ff00'>Oplossing:</font></b> Er mogen maximaal twee (sub)titelzinnen worden opgegeven (zinnen die beginnen met #). '</font>"
                                                f"<font size=5><b><font color='#335fff'>De presentatie is zonder deze zin(nen). </font>"
                                            )
                                            add_to_log(log_tekst)
                                            titel_error_weergeven = True
                                else:
                                    slidelijsten[0] = []
                                    slidelijsten[0].append(zin)

                                
                            elif zin.startswith("@"):
                                slidenummer = slidenummer + 1
                                zin = zin[1:] #haalt @ weg uit presentatie
                                slidelijsten[slidenummer] = [zin] 

                            elif zin.startswith("!"):

                                        
                                zin = zin[1:]  # haalt ! weg
                                # BRON: ChatGPT
                                # PROMPT: Leg uit waarom achtergrondkleur = "RGBColor(53, 75, 32) niet werkt"

                                # Uitleg: Het werkte niet omdat bij het maken van de slides het programma een string terugkreeg
                                # Dat begreep het programma niet omdat het losse kleurwaardes verwachtte (ipv R, G, B kreeg het RGB, -, - binnen)
                                
                                # Dit splits de gegeven kleurencodes in de R, G en de B. Benoemt het variabel en controleert of de kleurcode niet foutief is. 
                                # Als foutief achtergrondkleur = wit

                                try:
                                    # alleen achtergrondkleur selecteren als er geen errors waren bij kleurselecties om onverwachte uitkomsten stoppen.
                                    if kleur_error_woordgeving_weergeven == False or kleur_gekozen_error_weergeven == False:
                                        rgb_tuple = tuple(map(int, zin.strip("() ").split(",")))
                                        achtergrondkleur.append = RGBColor(*rgb_tuple)
                                        kleur_gekozen = True
                                except Exception as e:
                                    if kleur_error_woordgeving_weergeven == False:
                                            log_tekst = (
                                                f"<font size=5><b><font color='#ff0000'>Error: kon achtergrondkleur niet bepalen:</font></b> <i>'{zin}'</i> </font>"
                                                f"<font size=5><b><font color='#11ff00'>Oplossing:</font></b> Formateer de achtergrondkleur als '!(R, G, B).'. Correct voorbeeld: '!(122, 53, 36).' Getal mag 1 t/m 255 zijn. '</font>"
                                                f"<font size=5><b><font color='#335fff'>De achtergrond van de presentatie is wit gebleven. </font>"
                                            )
                                            add_to_log(log_tekst)
                                            kleur_error_woordgeving_weergeven = True
                                    achtergrondkleur = RGBColor(255, 255, 255) # Geen goede kleur dan witte achtergrond

                                if kleur_gekozen == False: # is errormessage al verschenen?
                                    if kleur_gekozen_error_weergeven == False:  #is er al een kleur gekozen
                                        log_tekst = (
                                            f"<font size=5><b><font color='#ff0000'>Error: Achtergrondkleur al gegeven:</font></b> <i>'{zin}'</i> </font>"
                                            f"<font size=5><b><font color='#11ff00'>Oplossing:</font></b> Er mag maximaal een achtergrondkleur worden opgegeven '</font>"
                                            f"<font size=5><b><font color='#335fff'>De achtergrondkleur van de presentatie is de eerste gegeven kleur geworden. </font>"
                                        )
                                        add_to_log(log_tekst)
                                        kleur_gekozen_error_weergeven = True

                            elif zin.startswith(">"):
                                zin = zin[1:]
                                onderdelen = [deel.strip() for deel in zin.split(",")]

                                bestandsnaam = onderdelen[0]
                                x = float(onderdelen[1]) if len(onderdelen) > 1 else 1.0 # Haalt x, y, breedte en hoogte op uit gegeven tekstbestand en Controleert of x is gegeven anders is waarde 1.0
                                y = float(onderdelen[2]) if len(onderdelen) > 2 else 1.0    
                                breedte = float(onderdelen[3]) if len(onderdelen) > 3 else None
                                hoogte = float(onderdelen[4]) if len(onderdelen) > 4 else 3.0

                                afbeelding_gevonden = False  # bijhouden of afbeelding is gevonden

                                afbeeldingen_map = os.path.join(os.path.dirname(__file__), "afbeeldingen") #Zoek het bestand in het afbeeldingenmapje
                                for bestand in os.listdir(afbeeldingen_map):
                                    if bestandsnaam.lower() in bestand.lower():
                                        slide_afbeelding = os.path.join(afbeeldingen_map, bestand)
                                        print(f"Gevonden bestand: {slide_afbeelding}")

                                        # alle gevens voor afbeelding op een slide
                                        afbeeldingen_per_slide[slidenummer] = {
                                            "pad": slide_afbeelding,
                                            "x": x,
                                            "y": y,
                                            "breedte": breedte,
                                            "hoogte": hoogte
                                        }

                                        afbeelding_gevonden = True
                                        break  # Afbeelding gevonden dus loopje stoppen
                                    
                                if not afbeelding_gevonden:
                                    log_tekst = (
                                            f"<font size=5><b><font color='#ff0000'>Error: Kon {bestandsnaam} niet vinden:</font></b> <i>'{zin}'</i> </font>"
                                            f"<font size=5><b><font color='#11ff00'>Oplossing:</font></b> Geef een correcte afbeelding op uit het afbeeldingmapje. '</font>"
                                            f"<font size=5><b><font color='#335fff'>De afbeelding is niet weergeven. </font>"
                                        )
                                    add_to_log(log_tekst)

                            else:
                                if slidenummer not in slidelijsten:
                                    #als de eerste zin(nen) geen @ of # hebben
                                    slidenummer = 1
                                    slidelijsten[slidenummer] = []
                                slidelijsten[slidenummer].append(zin)

                        prs = Presentation()

                        def titelslideophalen():
                            if 0 not in slidelijsten:
                                return  # geen titelslide

                            slide_layout = prs.slide_layouts[0]  # Titel-layout
                            slide = prs.slides.add_slide(slide_layout)
                            title = slide.shapes.title
                            subtitle = slide.placeholders[1]

                            background = slide.background
                            fill = background.fill
                            fill.solid()
                            #fill.fore_color.rgb = achtergrondkleur[0]

                            title.text = slidelijsten[0][0]
                            subtitle.text = "\n".join(slidelijsten[0][1:]).strip()

                        # Maak de inhoudslides
                        def inhoudslides():
                            for nummer, zinnen in slidelijsten.items():

                                if nummer == 0: #sla de gegevens van de titelslide over
                                    continue

                                #maakt slide
                                slide_layout = prs.slide_layouts[1]  # De layout (van de library) voor "Titel en inhoud"
                                slide = prs.slides.add_slide(slide_layout)

                                background = slide.background
                                fill = background.fill
                                fill.solid()
                                #for slidekleur in achtergrondkleur:
                                #    print(slidekleur)
                                #    print(achtergrondkleur)
                                #    #fill.fore_color.rgb = slidekleur
                                fill.fore_color.rgb = achtergrondkleur

                                #defining variables
                                title = slide.shapes.title
                                content = slide.placeholders[1]
                                text_frame = content.text_frame

                                # Afbeelding toevoegen
                                if nummer in afbeeldingen_per_slide:
                                    info = afbeeldingen_per_slide[nummer] # Eerder opgehaalde informatie wordt nu bij de correcte slide en afbeelding gezet voor in presentatie
                                    pad = info["pad"]
                                    x = Inches(info["x"]) #als x en y niet aanwezig is wordt het automatisch 1.0
                                    y = Inches(info["y"])
                                    breedte = Inches(info["breedte"]) if info["breedte"] else None #Als hoogte en breedte niet aanwezig is wordt het automatisch bepaalt
                                    hoogte = Inches(info["hoogte"]) if info["hoogte"] else None

                                    slide.shapes.add_picture(pad, x, y, width=breedte, height=hoogte)
                                    
                                for zin in zinnen[1:]:
                                    p_text = text_frame.add_paragraph()

                                    '''
                                    SOURCE: ChatGPT
                                    PROMPT: how do i get the characters between the brackets
                                    
                                    re import and how it works in this case
                                    '''

                                    for between_brackets in re.split(r'(\[.*?\])', zin):    #re is a weird import, unable to explain ¯\_(ツ)_/¯
                                        run = p_text.add_run()
                                        if between_brackets.startswith("[") and between_brackets.endswith("]"):
                                            run.text = between_brackets[1:-1]  # removes the brackets
                                            run.font.bold = True
                                        else:                                   #will just print the text if tehre is no closing bracket
                                            run.text = between_brackets

                                title = slide.shapes.title
                                p_title = title.text_frame.paragraphs[0]
                                p_title.clear()

                                for between_brackets in re.split(r'(\[.*?\])', zinnen[0]): #copying code for titles
                                    run = p_title.add_run()
                                    if between_brackets.startswith("[") and between_brackets.endswith("]"):
                                        run.text = between_brackets[1:-1]
                                        run.font.bold = True
                                    else:
                                        run.text = between_brackets


                        def bestandopslaan():
                            PPTXnaam = bestandnaamvak.get_text()
                            if PPTXnaam == "":
                                PPTXnaam = "Text-To-Slide"

                            downloads = os.path.join(os.path.expanduser("~"), "Downloads")
                            PPTXbestand = os.path.join(downloads, f"{PPTXnaam}.pptx")
                            bestandnaam = f"{PPTXnaam}.pptx"

                            achtergetal = 0
                            while os.path.exists(PPTXbestand):
                                achtergetal += 1
                                PPTXbestand = os.path.join(downloads, f"{PPTXnaam}({achtergetal}).pptx")
                                bestandnaam = f"{PPTXnaam}({achtergetal}).pptx"
                            prs.save(PPTXbestand)

                            log_tekst = (
                                f"<font size=5><b><font color='#FFFFFF'>{bestandnaam}</font></b> opgeslagen op de volgende locatie: <b><font color='#FFFFFF'>{PPTXbestand}</font></b></font>"
                                f"<font size=5><b>--------------------------------------------------------------------------</font></b>"
                            )
                            add_to_log(log_tekst)

                            print(f"Presentatie opgeslagen als: {PPTXbestand}")

                        titelslideophalen()
                        inhoudslides()
                        bestandopslaan()
                        titel_error_weergeven, kleur_error_woordgeving_weergeven, kleur_gekozen_error_weergeven, kleur_gekozen = False, False, False, False
                    
                    print("Gekozen bestand:", bestand)      
                    slidelijsten = bestandlezen(bestand)


                    
                else:
                    print("Geen bestand gekozen")

        manager.process_events(event)

    manager.update(time_delta)

    window_surface.blit(background, (0, 0))
    manager.draw_ui(window_surface)

    pygame.display.update()
