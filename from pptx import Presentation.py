from pptx import Presentation

# Maak een nieuwe presentatie
prs = Presentation()

slide_layout = prs.slide_layouts[0]  # Titel slide

slide = prs.slides.add_slide(slide_layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Titel  slide 1"
subtitle.text = "Ik maak graag pannenkoeken."

slide_layout2 = prs.slide_layouts[1]  # Titel + inhoud
slide2 = prs.slides.add_slide(slide_layout2)

title2 = slide2.shapes.title
content = slide2.placeholders[1]

title2.text = "Tweede slide"
content.text = "Eerste punnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnpunnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnntpunnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnntnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnnt\nTweede punt\nDerde punt"

# Opslaan als pptx bestand
prs.save("voorbeeld.pptx")
